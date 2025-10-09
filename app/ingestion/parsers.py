from __future__ import annotations

import hashlib
import io
import logging
import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional

from .chunking import TextBlock, chunk_blocks, chunk_text
from .types import ChunkPayload, ParsedDocument
from .utils import compute_checksum


logger = logging.getLogger(__name__)


@dataclass
class ParserContext:
    path: Path
    media_type: str


class DocumentParser:
    media_types: List[str] = []
    extensions: List[str] = []

    def matches(self, ctx: ParserContext) -> bool:
        if self.media_types and ctx.media_type in self.media_types:
            return True
        if self.extensions and ctx.path.suffix.lower() in self.extensions:
            return True
        return False

    def parse(self, ctx: ParserContext) -> ParsedDocument:
        raise NotImplementedError


class TextParser(DocumentParser):
    extensions = [".txt", ".md"]
    media_types = ["text/plain", "text/markdown"]

    def parse(self, ctx: ParserContext) -> ParsedDocument:
        text = ctx.path.read_text(encoding="utf-8")
        chunks = chunk_text(text)
        if not chunks:
            chunks = [ChunkPayload(ordering=0, text=text)]
        return ParsedDocument(
            display_name=ctx.path.name,
            media_type=ctx.media_type,
            checksum=compute_checksum(ctx.path),
            chunks=chunks,
            original_path=ctx.path,
        )


class PDFParser(DocumentParser):
    extensions = [".pdf"]
    media_types = ["application/pdf"]

    def parse(self, ctx: ParserContext) -> ParsedDocument:
        import pdfplumber

        chunks: List[ChunkPayload] = []
        with pdfplumber.open(ctx.path) as pdf:
            for idx, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if not text.strip():
                    continue
                page_chunks = chunk_text(
                    text.strip(),
                    page_label=str(page.page_number),
                )
                if not page_chunks:
                    continue
                for chunk in page_chunks:
                    chunk.ordering = len(chunks)
                    chunks.append(chunk)
        return ParsedDocument(
            display_name=ctx.path.name,
            media_type=ctx.media_type,
            checksum=compute_checksum(ctx.path),
            chunks=chunks,
            original_path=ctx.path,
        )


class DocxParser(DocumentParser):
    extensions = [".docx"]
    media_types = [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ]

    def parse(self, ctx: ParserContext) -> ParsedDocument:
        import docx
        document = docx.Document(str(ctx.path))
        text_blocks = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        if text_blocks:
            block_payloads = [TextBlock(text=block) for block in text_blocks]
            chunks = chunk_blocks(block_payloads)
        else:
            chunks = []
        if not chunks:
            full_text = "\n".join(text_blocks) if text_blocks else ""
            chunks = chunk_text(full_text or "") or [ChunkPayload(ordering=0, text=full_text)]
        for idx, chunk in enumerate(chunks):
            chunk.ordering = idx
        return ParsedDocument(
            display_name=ctx.path.name,
            media_type=ctx.media_type,
            checksum=compute_checksum(ctx.path),
            chunks=chunks,
            original_path=ctx.path,
        )


class PptxParser(DocumentParser):
    extensions = [".pptx"]
    media_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]

    def parse(self, ctx: ParserContext) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        try:
            import pytesseract  # type: ignore[import]
        except ImportError as exc:  # pragma: no cover - dependency missing at runtime
            raise RuntimeError(
                "pytesseract is required for PPTX image OCR. Install it via pip and ensure the"
                " Tesseract OCR engine is available on the system PATH."
            ) from exc

        try:
            from PIL import Image  # type: ignore[import]
        except ImportError as exc:  # pragma: no cover - dependency missing at runtime
            raise RuntimeError(
                "Pillow is required for PPTX image OCR. Install it via pip before ingesting slides"
                " that contain images."
            ) from exc

        presentation = Presentation(str(ctx.path))
        chunks: List[ChunkPayload] = []
        seen_hashes: set[str] = set()
        image_counter = 0
        # Collect all slide texts and OCR results as (slide_idx, text, provenance_type, extra_metadata)
        slide_texts = []
        for idx, slide in enumerate(presentation.slides):
            lines: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(getattr(shape, "text", "")).strip()
                    if text:
                        lines.append(text)
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    image = getattr(shape, "image", None)
                    if image is None:
                        continue
                    blob = image.blob
                    digest = hashlib.sha256(blob).hexdigest()
                    if digest in seen_hashes:
                        continue
                    seen_hashes.add(digest)
                    image_counter += 1
                    try:
                        with Image.open(io.BytesIO(blob)) as img:
                            if img.mode in ('RGBA', 'LA', 'P'):
                                img = img.convert('RGB')
                            ocr_text = pytesseract.image_to_string(img)
                    except pytesseract.TesseractNotFoundError as exc:
                        raise RuntimeError(
                            "Tesseract OCR executable not found. Install it locally and ensure it is"
                            " accessible on the system PATH before ingesting PPTX images."
                        ) from exc
                    except (TypeError, OSError) as exc:
                        logger.warning(
                            "Skipping unsupported image format in slide %s image %s of %s: %s",
                            idx + 1,
                            image_counter,
                            ctx.path,
                            exc,
                        )
                        continue
                    except Exception:  # noqa: BLE001
                        logger.exception(
                            "Failed to run OCR for slide %s image %s in %s",
                            idx + 1,
                            image_counter,
                            ctx.path,
                        )
                        continue

                    cleaned = (ocr_text or "").strip()
                    if cleaned:
                        slide_texts.append((idx + 1, cleaned, "image_ocr", {
                            "image_index": image_counter,
                            "image_hash": digest,
                        }))
            if lines:
                block_text = "\n".join(lines)
                slide_texts.append((idx + 1, block_text, "slide_text", {}))

        # Now batch across slides until token limit is hit
        from .chunking import chunk_blocks, TextBlock, DEFAULT_MAX_TOKENS, DEFAULT_OVERLAP_TOKENS, DEFAULT_ENCODING
        current_texts = []
        current_slide_indices = []
        current_types = []
        current_metadata = []
        encoder = None
        for slide_idx, text, prov_type, extra_meta in slide_texts:
            if encoder is None:
                import tiktoken
                encoder = tiktoken.get_encoding(DEFAULT_ENCODING)
            # If adding this text would exceed the token limit, flush current chunk
            tokens_in_text = len(encoder.encode(text))
            tokens_in_current = len(encoder.encode("\n".join(current_texts))) if current_texts else 0
            if current_texts and (tokens_in_current + tokens_in_text > DEFAULT_MAX_TOKENS):
                chunk_text_combined = "\n".join(current_texts)
                chunk_blocks_result = chunk_blocks(
                    [TextBlock(text=chunk_text_combined)],
                    max_tokens=DEFAULT_MAX_TOKENS,
                    overlap_tokens=DEFAULT_OVERLAP_TOKENS,
                    encoding_name=DEFAULT_ENCODING,
                )
                slide_range = f"{current_slide_indices[0]}-{current_slide_indices[-1]}" if len(current_slide_indices) > 1 else str(current_slide_indices[0])
                for chunk in chunk_blocks_result:
                    chunk.ordering = len(chunks)
                    chunk.metadata = {
                        "source": ",".join(set(current_types)),
                        "slide_range": slide_range,
                        "slide_indices": list(current_slide_indices),
                        **(current_metadata[0] if current_metadata else {}),
                    }
                    chunks.append(chunk)
                current_texts = []
                current_slide_indices = []
                current_types = []
                current_metadata = []
            current_texts.append(text)
            current_slide_indices.append(slide_idx)
            current_types.append(prov_type)
            current_metadata.append(extra_meta)
        # Flush any remaining text
        if current_texts:
            chunk_text_combined = "\n".join(current_texts)
            chunk_blocks_result = chunk_blocks(
                [TextBlock(text=chunk_text_combined)],
                max_tokens=DEFAULT_MAX_TOKENS,
                overlap_tokens=DEFAULT_OVERLAP_TOKENS,
                encoding_name=DEFAULT_ENCODING,
            )
            slide_range = f"{current_slide_indices[0]}-{current_slide_indices[-1]}" if len(current_slide_indices) > 1 else str(current_slide_indices[0])
            for chunk in chunk_blocks_result:
                chunk.ordering = len(chunks)
                chunk.metadata = {
                    "source": ",".join(set(current_types)),
                    "slide_range": slide_range,
                    "slide_indices": list(current_slide_indices),
                    **(current_metadata[0] if current_metadata else {}),
                }
                chunks.append(chunk)
        return ParsedDocument(
            display_name=ctx.path.name,
            media_type=ctx.media_type,
            checksum=compute_checksum(ctx.path),
            chunks=chunks,
            original_path=ctx.path,
        )


class ParserRegistry:
    def __init__(self, parsers: Optional[List[DocumentParser]] = None):
        self.parsers = parsers or [PDFParser(), DocxParser(), PptxParser(), TextParser()]

    def detect_media_type(self, path: Path) -> str:
        guessed, _ = mimetypes.guess_type(path.resolve().as_uri())
        return guessed or "application/octet-stream"

    def select_parser(self, path: Path) -> DocumentParser:
        media_type = self.detect_media_type(path)
        ctx = ParserContext(path=path, media_type=media_type)
        for parser in self.parsers:
            if parser.matches(ctx):
                return parser
        raise ValueError(f"Unsupported file type: {path.suffix}")

    def parse(self, path: Path) -> ParsedDocument:
        media_type = self.detect_media_type(path)
        ctx = ParserContext(path=path, media_type=media_type)
        parser = self.select_parser(path)
        return parser.parse(ctx)


DEFAULT_REGISTRY = ParserRegistry()


def parse_document(path: Path) -> ParsedDocument:
    return DEFAULT_REGISTRY.parse(path)
